"""
Extract product images + labels from a PowerPoint line-sheet deck.

Input: a .pptx where products are laid out as image shapes with nearby text boxes
containing product name, style-color code, wholesale price, and retail price.

Outputs:
- extracted_product_data.xlsx
- extracted_product_data.csv
- images_raw/       original embedded picture blobs
- images_clean/     PNG images cropped to visible product area where possible

Install locally:
    pip install python-pptx pillow pandas openpyxl

Run:
    python ppt_product_extract.py sample.pptx --out output_folder
"""

from __future__ import annotations

import argparse
import csv
import hashlib
import io
import re
from dataclasses import dataclass, asdict
from pathlib import Path
from typing import Any, Iterable

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image, ImageChops

try:
    import pandas as pd
except Exception:  # keep CSV output possible even without pandas
    pd = None

STYLE_RE = re.compile(r"\b[A-Z]{1,3}\d{4}-[A-Z0-9]+(?:/[A-Z0-9]+)*\b")
PRICE_RE = re.compile(
    r"\$\s*([0-9]+(?:\.[0-9]{1,2})?)\s*WS\s*\|\s*\$\s*([0-9]+(?:\.[0-9]{1,2})?)\s*RETAIL",
    re.IGNORECASE,
)

HEADER_WORDS = {
    "NIKE GLOBAL SPORTS APPAREL",
    "KEEP IT TIGHT.",
    "KEEP IT TIGHT",
    "WOMEN’S APPAREL",
    "WOMEN'S APPAREL",
    "MEN’S APPAREL",
    "MEN'S APPAREL",
}


@dataclass
class TextItem:
    slide_number: int
    slide_title: str
    shape_name: str
    raw_text: str
    product_name: str
    style_color_code: str
    style_code: str
    color_code: str
    wholesale_price: float | None
    retail_price: float | None
    text_left: int
    text_top: int
    text_width: int
    text_height: int


@dataclass
class PictureItem:
    slide_number: int
    image_index: int
    shape_name: str
    left: int
    top: int
    width: int
    height: int
    original_filename: str
    clean_filename: str | None
    image_ext: str
    image_sha1: str
    average_hex: str | None
    dominant_hex_1: str | None
    dominant_hex_2: str | None
    dominant_hex_3: str | None
    visual_hash: str | None


def iter_shapes(shapes: Iterable[Any]) -> Iterable[Any]:
    """Yield shapes recursively. Useful if future decks use grouped objects."""
    for shape in shapes:
        if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.GROUP:
            yield from iter_shapes(shape.shapes)
        else:
            yield shape


def normalize_text(text: str) -> str:
    text = text.replace("\u00a0", " ")
    return "\n".join(line.strip() for line in text.splitlines() if line.strip())


def parse_product_text(slide_number: int, slide_title: str, shape: Any) -> TextItem | None:
    raw = normalize_text(getattr(shape, "text", "") or "")
    if not raw:
        return None

    style_match = STYLE_RE.search(raw)
    price_match = PRICE_RE.search(raw)
    if not style_match or not price_match:
        return None

    style_color = style_match.group(0).strip()
    style_code, color_code = style_color.split("-", 1)

    # Product name is the text before the style-color code, with line breaks collapsed.
    before_style = raw[: style_match.start()]
    product_name = " ".join(line.strip() for line in before_style.splitlines() if line.strip())
    product_name = re.sub(r"\s+", " ", product_name).strip()

    ws_price = float(price_match.group(1)) if price_match else None
    retail_price = float(price_match.group(2)) if price_match else None

    return TextItem(
        slide_number=slide_number,
        slide_title=slide_title,
        shape_name=shape.name,
        raw_text=raw,
        product_name=product_name,
        style_color_code=style_color,
        style_code=style_code,
        color_code=color_code,
        wholesale_price=ws_price,
        retail_price=retail_price,
        text_left=int(shape.left),
        text_top=int(shape.top),
        text_width=int(shape.width),
        text_height=int(shape.height),
    )


def get_slide_title(slide: Any) -> str:
    candidates: list[str] = []
    for shape in iter_shapes(slide.shapes):
        if not hasattr(shape, "text"):
            continue
        text = normalize_text(shape.text)
        if not text:
            continue
        first_line = text.splitlines()[0].strip()
        if first_line.upper() in HEADER_WORDS or "APPAREL" in first_line.upper():
            candidates.append(first_line)
    return candidates[0] if candidates else ""


def crop_visible(im: Image.Image) -> Image.Image:
    """Crop transparent or near-white border around product images."""
    rgba = im.convert("RGBA")

    # Prefer alpha-based crop when transparent background exists.
    alpha = rgba.getchannel("A")
    alpha_bbox = alpha.point(lambda p: 255 if p > 10 else 0).getbbox()
    if alpha_bbox:
        rgba = rgba.crop(alpha_bbox)

    # Then remove almost-white borders.
    rgb = rgba.convert("RGB")
    bg = Image.new("RGB", rgb.size, (255, 255, 255))
    diff = ImageChops.difference(rgb, bg)
    # amplify subtle differences from off-white background
    diff = ImageChops.add(diff, diff, 2.0, -20)
    bbox = diff.getbbox()
    if bbox:
        rgba = rgba.crop(bbox)
    return rgba


def hex_from_rgb(rgb: tuple[int, int, int]) -> str:
    return "#%02X%02X%02X" % rgb


def image_color_features(im: Image.Image) -> tuple[str | None, list[str | None], str | None]:
    """Return average color, top 3 dominant colors, and an 8x8 average hash."""
    rgba = im.convert("RGBA").resize((80, 80))
    pixels = list(rgba.getdata())
    visible = [p for p in pixels if p[3] > 20 and not (p[0] > 245 and p[1] > 245 and p[2] > 245)]
    if not visible:
        return None, [None, None, None], None

    avg = tuple(int(sum(p[i] for p in visible) / len(visible)) for i in range(3))
    avg_hex = hex_from_rgb(avg)

    # Dominant colors using PIL adaptive palette on non-transparent crop composited on white.
    rgb = im.convert("RGBA")
    bg = Image.new("RGBA", rgb.size, (255, 255, 255, 255))
    comp = Image.alpha_composite(bg, rgb).convert("RGB").resize((80, 80))
    pal = comp.quantize(colors=5, method=Image.Quantize.MEDIANCUT)
    palette = pal.getpalette() or []
    counts = pal.getcolors() or []
    ranked: list[str] = []
    for count, idx in sorted(counts, reverse=True):
        r, g, b = palette[idx * 3 : idx * 3 + 3]
        # skip near-white background
        if r > 245 and g > 245 and b > 245:
            continue
        hx = hex_from_rgb((r, g, b))
        if hx not in ranked:
            ranked.append(hx)
        if len(ranked) == 3:
            break
    ranked += [None] * (3 - len(ranked))

    gray = im.convert("L").resize((8, 8))
    vals = list(gray.getdata())
    threshold = sum(vals) / len(vals)
    bits = "".join("1" if v >= threshold else "0" for v in vals)
    visual_hash = f"{int(bits, 2):016x}"
    return avg_hex, ranked[:3], visual_hash


def save_picture(shape: Any, slide_number: int, image_index: int, out_dir: Path) -> PictureItem | None:
    if getattr(shape, "shape_type", None) != MSO_SHAPE_TYPE.PICTURE:
        return None

    img = shape.image
    ext = (img.ext or "bin").lower()
    blob = img.blob
    sha1 = hashlib.sha1(blob).hexdigest()

    raw_dir = out_dir / "images_raw"
    clean_dir = out_dir / "images_clean"
    raw_dir.mkdir(parents=True, exist_ok=True)
    clean_dir.mkdir(parents=True, exist_ok=True)

    base = f"slide_{slide_number:03d}_img_{image_index:03d}_{shape.name.replace(' ', '_')}"
    raw_filename = f"images_raw/{base}.{ext}"
    raw_path = out_dir / raw_filename
    raw_path.write_bytes(blob)

    clean_filename = None
    avg_hex = None
    doms: list[str | None] = [None, None, None]
    visual_hash = None

    try:
        im = Image.open(io.BytesIO(blob))
        clean = crop_visible(im)
        clean_filename = f"images_clean/{base}.png"
        clean_path = out_dir / clean_filename
        clean.save(clean_path)
        avg_hex, doms, visual_hash = image_color_features(clean)
    except Exception:
        # Some embedded image formats may not be readable by PIL; keep raw file at minimum.
        pass

    return PictureItem(
        slide_number=slide_number,
        image_index=image_index,
        shape_name=shape.name,
        left=int(shape.left),
        top=int(shape.top),
        width=int(shape.width),
        height=int(shape.height),
        original_filename=raw_filename,
        clean_filename=clean_filename,
        image_ext=ext,
        image_sha1=sha1,
        average_hex=avg_hex,
        dominant_hex_1=doms[0],
        dominant_hex_2=doms[1],
        dominant_hex_3=doms[2],
        visual_hash=visual_hash,
    )


def overlap_1d(a0: int, a1: int, b0: int, b1: int) -> int:
    return max(0, min(a1, b1) - max(a0, b0))


def picture_match_score(item: TextItem, pic: PictureItem, slide_w: int, slide_h: int) -> tuple[float, str]:
    """Score a text/image pair. Lower is better."""
    tx0, tx1 = item.text_left, item.text_left + item.text_width
    ty0 = item.text_top
    tcx = item.text_left + item.text_width / 2

    px0, px1 = pic.left, pic.left + pic.width
    py1 = pic.top + pic.height
    pcx = pic.left + pic.width / 2
    pcy = pic.top + pic.height / 2

    xdiff = abs(tcx - pcx) / slide_w
    vertical_gap = (ty0 - py1) / slide_h  # positive when text is below picture
    y_penalty = 0.0 if pcy < ty0 else 0.9
    gap_penalty = abs(vertical_gap) if vertical_gap >= -0.08 else abs(vertical_gap) + 0.4
    x_overlap = overlap_1d(tx0, tx1, px0, px1) / max(1, min(item.text_width, pic.width))
    overlap_bonus = 0.30 * x_overlap
    size_penalty = 0.15 if pic.width < slide_w * 0.035 or pic.height < slide_h * 0.035 else 0.0

    score = 3.0 * xdiff + gap_penalty + y_penalty + size_penalty - overlap_bonus
    reason = f"xdiff={xdiff:.3f};vgap={vertical_gap:.3f};xoverlap={x_overlap:.2f}"
    return score, reason


def match_pictures_global(items: list[TextItem], pictures: list[PictureItem], slide_w: int, slide_h: int) -> dict[int, tuple[int | None, float | None, str]]:
    """Assign pictures to labels using global greedy matching by best pair score.

    This avoids early text boxes consuming images that are better matches for later text boxes.
    """
    pairs: list[tuple[float, int, int, str]] = []
    for item_idx, item in enumerate(items):
        for pic_idx, pic in enumerate(pictures):
            score, reason = picture_match_score(item, pic, slide_w, slide_h)
            pairs.append((score, item_idx, pic_idx, reason))

    assignments: dict[int, tuple[int | None, float | None, str]] = {}
    used_items: set[int] = set()
    used_pics: set[int] = set()

    for score, item_idx, pic_idx, reason in sorted(pairs, key=lambda x: x[0]):
        if item_idx in used_items or pic_idx in used_pics:
            continue
        used_items.add(item_idx)
        used_pics.add(pic_idx)
        confidence = max(0.0, min(1.0, 1.0 - score))
        assignments[item_idx] = (pic_idx, round(confidence, 3), reason)
        if len(used_items) == len(items) or len(used_pics) == len(pictures):
            break

    for item_idx in range(len(items)):
        assignments.setdefault(item_idx, (None, None, "no candidate"))
    return assignments


def extract_pptx(pptx_path: Path, out_dir: Path) -> list[dict[str, Any]]:
    prs = Presentation(str(pptx_path))
    out_dir.mkdir(parents=True, exist_ok=True)

    rows: list[dict[str, Any]] = []
    all_pictures_count = 0

    for slide_idx, slide in enumerate(prs.slides, start=1):
        slide_title = get_slide_title(slide)
        text_items: list[TextItem] = []
        pictures: list[PictureItem] = []

        image_index = 1
        for shape in iter_shapes(slide.shapes):
            parsed = parse_product_text(slide_idx, slide_title, shape) if hasattr(shape, "text") else None
            if parsed:
                text_items.append(parsed)
            if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE:
                pic = save_picture(shape, slide_idx, image_index, out_dir)
                if pic:
                    pictures.append(pic)
                    all_pictures_count += 1
                    image_index += 1

        sorted_items = sorted(text_items, key=lambda x: (x.text_top, x.text_left))
        assignments = match_pictures_global(sorted_items, pictures, int(prs.slide_width), int(prs.slide_height))
        for item_no, item in enumerate(sorted_items, start=1):
            match_idx, confidence, reason = assignments.get(item_no - 1, (None, None, "no candidate"))
            pic = pictures[match_idx] if match_idx is not None else None

            row = asdict(item)
            row.update(
                {
                    "item_number_on_slide": item_no,
                    "matched_image_confidence": confidence,
                    "matched_image_reason": reason,
                    "image_shape_name": pic.shape_name if pic else None,
                    "image_file": pic.clean_filename if pic and pic.clean_filename else None,
                    "raw_image_file": pic.original_filename if pic else None,
                    "image_ext": pic.image_ext if pic else None,
                    "image_sha1": pic.image_sha1 if pic else None,
                    "image_left": pic.left if pic else None,
                    "image_top": pic.top if pic else None,
                    "image_width": pic.width if pic else None,
                    "image_height": pic.height if pic else None,
                    "average_hex": pic.average_hex if pic else None,
                    "dominant_hex_1": pic.dominant_hex_1 if pic else None,
                    "dominant_hex_2": pic.dominant_hex_2 if pic else None,
                    "dominant_hex_3": pic.dominant_hex_3 if pic else None,
                    "visual_hash": pic.visual_hash if pic else None,
                    "price_multiple_retail_over_ws": round(item.retail_price / item.wholesale_price, 4)
                    if item.retail_price and item.wholesale_price
                    else None,
                }
            )
            rows.append(row)

    csv_path = out_dir / "extracted_product_data.csv"
    if rows:
        with csv_path.open("w", newline="", encoding="utf-8-sig") as f:
            writer = csv.DictWriter(f, fieldnames=list(rows[0].keys()))
            writer.writeheader()
            writer.writerows(rows)

    if pd is not None and rows:
        xlsx_path = out_dir / "extracted_product_data.xlsx"
        df = pd.DataFrame(rows)
        summary = pd.DataFrame(
            [
                {"metric": "pptx_file", "value": str(pptx_path)},
                {"metric": "slides_scanned", "value": len(prs.slides)},
                {"metric": "product_text_rows_extracted", "value": len(rows)},
                {"metric": "pictures_extracted", "value": all_pictures_count},
                {"metric": "note", "value": "Image mapping is geometry-based; review low-confidence rows."},
            ]
        )
        with pd.ExcelWriter(xlsx_path, engine="openpyxl") as writer:
            df.to_excel(writer, sheet_name="Products", index=False)
            summary.to_excel(writer, sheet_name="Summary", index=False)

    return rows


def main() -> None:
    parser = argparse.ArgumentParser(description="Extract product images and labels from a PowerPoint line-sheet.")
    parser.add_argument("pptx", type=Path, help="Path to input .pptx")
    parser.add_argument("--out", type=Path, default=Path("ppt_extract_output"), help="Output folder")
    args = parser.parse_args()

    rows = extract_pptx(args.pptx, args.out)
    print(f"Extracted {len(rows)} product rows into: {args.out}")
    print(f"CSV:  {args.out / 'extracted_product_data.csv'}")
    if pd is not None:
        print(f"XLSX: {args.out / 'extracted_product_data.xlsx'}")


if __name__ == "__main__":
    main()
