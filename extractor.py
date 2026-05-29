from __future__ import annotations

import io
import re
import zipfile
from pathlib import Path
from typing import Iterable

import pandas as pd
from PIL import Image, UnidentifiedImageError
from pptx import Presentation


# -----------------------------
# Regex / cleaning rules
# -----------------------------

STYLE_RE = re.compile(r"\b(?=[A-Z0-9]{6}\b)(?=[A-Z0-9]*\d)[A-Z]{2}[A-Z0-9]{4}\b")
# Standard Nike style-color cards, including slash colors such as IR5946-010/676.
STANDARD_STYLE_COLOR_RE = re.compile(r"\b([A-Z]{2}[A-Z0-9]{4})-(\d{3})(?:/(\d{3}))?\b")
STYLE_COLOR_RE = STANDARD_STYLE_COLOR_RE
DETAIL_STYLE_LINE_RE = re.compile(
    r"^(?P<style>[A-Z0-9]{6})\s*\|\s*\$(?P<ws>\d+(?:\.\d{1,2})?)\s*WHOLESALE\s*\|\s*\$(?P<retail>\d+(?:\.\d{1,2})?)\s*RETAIL\s*\|\s*(?P<status>NEW|CARRYOVER)$",
    re.I,
)

# Special Nike accessories/bags style-color format, for example: N.100.3478.091
ACCESSORY_STYLE_COLOR_RE = re.compile(r"\bN\.\d{3}\.\d{4}\.\d{3}\b", re.I)
VALID_STANDARD_STYLE_COLOR_RE = re.compile(r"^[A-Z]{2}[A-Z0-9]{4}-\d{3}$")
VALID_ACCESSORY_STYLE_COLOR_RE = re.compile(r"^N\.\d{3}\.\d{4}\.\d{3}$", re.I)

STATUS_RE = re.compile(r"\b(NEW|CARRYOVER)\b", re.I)

# IMPORTANT: use word boundaries so WSH is not interpreted as "WS".
PRICE_CONTEXT_RE = re.compile(r"\b(RETAIL|WS|WHSE|WHOLESALE)\b", re.I)

STOP_RE = re.compile(
    r"^(Material|Material Content|Features and Benefits|Key Features|Overview|Product Details|Keep it|FOR INTERNAL USE|©|Product Images|"
    r"Performance Benefits|FIT:|RIDE:|TRACTION:|Product Description)",
    re.I,
)

# Lines that often appear near the style code but are NOT product names.
BAD_PRODUCT_NAME_RE = re.compile(
    r"^(?:"
    r"Dri-Fit|Therma-Fit|Aero-Fit(?: Technology)?|Storm-Fit(?: ADV)?|"
    r"Repel Technology|UV Protective|UV Protection|"
    r"Premium knit fabric|Premium knit jacquard fabric|Premium hand feel|"
    r"Standard Fit|Slim Fit|Loose Fit|Oversized Fit|Tight fit|Secure Fit|"
    r"Material:?|Key Features:?|Product Information|TBD|LAUNCH TBD|"
    r"SP\d+|SPSU\d+|"
    r"Stretch to provide full range of motion|Thumbholes|Snapback closure|"
    r"Stretchy, elastic closure|100% Cotton|100% Polyester|"
    r"Mid-depth|Max-depth|High depth|Low depth"
    r")$",
    re.I,
)

FEATURE_PRODUCT_NAME_RE = re.compile(
    r"(?i)^(?:"
    r".*\d+%\s*(?:polyester|cotton|nylon|spandex|elastane|rayon|modal).*|"
    r"(?:body|lining|material content)\s*:.*|"
    r"embroidered\s+(?:shield|swoosh|logo).*|"
    r".*\b(?:shield|swoosh|logo)\b.*(?:embroidery|embroidered|on\s+(?:sleeve|front|chest|back)).*|"
    r".*(?:embroidery|embroidered).*\b(?:shield|swoosh|logo)\b.*|"
    r".*\b(?:shield|swoosh|logo)\s+on\b.*|"
    r"top\s+stitched.*|"
    r".*\b(?:hood|3-piece hood)\s+construction\b.*|"
    r"utility\s+webbing.*|"
    r"strap-through\s+design.*|"
    r"faux-fur\s+lining.*|"
    r"premium$|"
    r"synth\s*etic$|"
    r"(?:drivers?|putters?|mallets?|headcovers?|head\s*covers?|bags?)$|"
    r".*\b(?:helps|protects|holds|features and benefits|key features)\b.*|"
    r"N\.\d{3}\.\d{4}\.\d{3}.*|"
    r".*\((?:left|right)\)\s*/?$"
    r")$"
)

HEADER_PRODUCT_NAME_RE = re.compile(
    r"(?i)^(?:"
    r"NIKE\s+GLOBAL\s+SPORTS\s+APPAREL.*|"
    r"KEEP\s+IT\s+TIGHT\.?|"
    r"SEASON\s*/\s*S\d+|"
    r"V\.\d+|"
    r"©.*NIKE.*"
    r")$"
)

# Strong title patterns for Nike product names.
PRODUCT_TITLE_RE = re.compile(
    r"^(?:M|W|B|G|K|U)\s+NK\b|"
    r"^NK\s+|"
    r"^K\s+NSW\b|"
    r"^W\s+NSW\b|"
    r"^NIKE\s+|"
    r"^Nike Brasilia\b|"
    r"^VICTORY\b|"
    r"^NEXT%\b|"
    r"^PEGASUS\b|"
    r"^TEMPO\b|"
    r"^AIR MAX\b|"
    r"^WOMEN['’]S\s+",
    re.I,
)

PRODUCT_KEYWORDS = {
    "POLO", "PANT", "SHORT", "JACKET", "JKT", "HOODIE", "VEST",
    "CAP", "SKIRT", "SKRT", "SHOE", "SOCK", "SOCKS", "BACKPACK",
    "BAG", "TOP", "TEE", "CREW", "DRESS", "VISOR", "BEANIE",
    "DUFFEL", "DRAWSTRING", "FOOTIE", "NO-SHOW", "ANKLE", "TRAINING",
}

CATEGORY_PATTERNS: list[tuple[str, re.Pattern]] = [
    ("Men's Apparel", re.compile(r"^(?:MEN'?S|MENS)\s+APPAREL$", re.I)),
    ("Women's Apparel", re.compile(r"^(?:WOMEN'?S|WOMENS)\s+APPAREL$", re.I)),
    ("Kids' Apparel", re.compile(r"^(?:KIDS?|YOUTH|BOYS?|GIRLS?)\s+APPAREL$", re.I)),
    ("Apparel", re.compile(r"^APPAREL$", re.I)),
    ("Bags", re.compile(r"^(?:BAGS?|BACKPACKS?|DUFFELS?|DRAWSTRING\s+BAGS?)$", re.I)),
    ("Gloves", re.compile(r"^GLOVES?$", re.I)),
    ("Headcovers", re.compile(r"^(?:HEAD\s*COVERS?|HEADCOVERS?)$", re.I)),
    ("Footwear", re.compile(r"^(?:FOOTWEAR|SHOES?)$", re.I)),
    ("Socks", re.compile(r"^SOCKS?$", re.I)),
    ("Caps & Headwear", re.compile(r"^(?:CAPS?|HATS?|HEADWEAR|VISORS?|BEANIES?)$", re.I)),
    ("Accessories", re.compile(r"^ACCESSORIES$", re.I)),
]


# -----------------------------
# Helper functions
# -----------------------------

def clean_text(value: str) -> str:
    value = str(value)
    value = value.replace("\xa0", " ")
    value = value.replace("–", "-").replace("—", "-")
    value = value.replace("’", "'")
    value = value.replace("–", "-").replace("—", "-")
    return re.sub(r"\s+", " ", value).strip()


def unique_preserve_order(values: list[str]) -> list[str]:
    result: list[str] = []
    for value in values:
        if value and value not in result:
            result.append(value)
    return result


def extract_style_color_pairs(text: str) -> list[tuple[str, str]]:
    text = clean_text(text)
    pairs: list[tuple[str, str]] = []
    for match in STYLE_COLOR_RE.finditer(text):
        style = match.group(1).upper()
        pairs.append((style, match.group(2).zfill(3)))
        if match.group(3):
            pairs.append((style, match.group(3).zfill(3)))
    return pairs


def expand_standard_style_color_tokens(text: str) -> list[dict]:
    """
    Expand tokens such as IR5946-010/676 into one row per color.
    """
    results: list[dict] = []
    for match in STYLE_COLOR_RE.finditer(clean_text(text)):
        style = match.group(1).upper()
        raw_code = match.group(0).upper().replace("â€“", "-").replace("â€”", "-")
        colors = [match.group(2).zfill(3)]
        if match.group(3):
            colors.append(match.group(3).zfill(3))

        for color in colors:
            results.append(
                {
                    "style_code": style,
                    "color_code": color,
                    "style_color_key": build_style_color_key(style, color),
                    "original_raw_code": raw_code,
                    "slash_color_expanded": len(colors) > 1,
                }
            )
    return results


def is_valid_style_color_key(value: str) -> bool:
    key = clean_text(value).upper()
    standard_match = VALID_STANDARD_STYLE_COLOR_RE.fullmatch(key)
    if standard_match:
        style = key.split("-", 1)[0]
        return any(ch.isdigit() for ch in style)

    return bool(
        VALID_ACCESSORY_STYLE_COLOR_RE.fullmatch(key)
    )


def parse_detail_style_line(line: str) -> dict | None:
    match = DETAIL_STYLE_LINE_RE.fullmatch(clean_text(line).upper())
    if not match:
        return None

    style = match.group("style").upper()
    if not any(ch.isdigit() for ch in style):
        return None

    return {
        "style": style,
        "wholesale": float(match.group("ws")),
        "retail": float(match.group("retail")),
        "status": match.group("status").upper(),
    }


def extract_colors_from_style_color_tokens(text: str) -> list[str]:
    return unique_preserve_order([color for _, color in extract_style_color_pairs(text)])


def detect_slide_category(lines: list[str]) -> str | None:
    """
    Detect section divider/category slides such as WOMEN'S APPAREL or HEADCOVERS.
    Product-heavy lines are ignored so product names do not accidentally reset category.
    """
    candidates = lines[:12]

    for line in candidates:
        line_clean = clean_text(line)
        if not line_clean:
            continue

        if STYLE_RE.search(line_clean) or ACCESSORY_STYLE_COLOR_RE.search(line_clean):
            continue

        if looks_like_price_context(line_clean):
            continue

        for category, pattern in CATEGORY_PATTERNS:
            if pattern.search(line_clean):
                return category

    return None


def iter_all_shapes(shapes):
    """Yield shapes recursively, including shapes inside groups."""
    for shape in shapes:
        yield shape
        if hasattr(shape, "shapes"):
            yield from iter_all_shapes(shape.shapes)


def get_slide_text_items(slide, slide_width: int | None = None, slide_height: int | None = None) -> list[dict]:
    """
    Extract text from slide shapes with their visual position.
    Sorting by top/left makes the extraction closer to what we see on the slide.
    """
    items: list[dict] = []

    for shape_index, shape in enumerate(iter_all_shapes(slide.shapes)):
        if not hasattr(shape, "text") or not shape.text:
            continue

        left = int(getattr(shape, "left", 0))
        top = int(getattr(shape, "top", 0))
        width = int(getattr(shape, "width", 0))
        height = int(getattr(shape, "height", 0))

        if slide_width is not None and slide_height is not None:
            completely_off_slide = (
                left + width < 0
                or top + height < 0
                or left > slide_width
                or top > slide_height
            )
            if completely_off_slide:
                continue

        text = shape.text.strip()
        if not text:
            continue

        for line_index, raw_line in enumerate(text.splitlines()):
            line = clean_text(raw_line)
            if not line:
                continue

            items.append(
                {
                    "shape_index": shape_index,
                    "line_index": line_index,
                    "text": line,
                    "left": left,
                    "top": top,
                    "width": width,
                    "height": height,
                }
            )

    # Visual order: top to bottom, then left to right
    return sorted(items, key=lambda x: (x["top"], x["left"], x["shape_index"]))


def get_slide_lines(slide, slide_width: int | None = None, slide_height: int | None = None) -> list[str]:
    return [item["text"] for item in get_slide_text_items(slide, slide_width, slide_height)]


def looks_like_price_context(text: str) -> bool:
    return bool(PRICE_CONTEXT_RE.search(clean_text(text)))


def is_new_product_start(line: str, lookahead: str = "") -> bool:
    """
    A real product detail line must contain a style code and WHOLESALE/RETAIL context.

    This intentionally ignores overview/product-grid lines like:
        WS: IR4565
        $140/70
    because those are not detailed product blocks and usually do not contain colors.
    """
    line_clean = clean_text(line)
    if parse_detail_style_line(line_clean):
        return True

    # Ignore overview/grid lines.
    if re.match(r"^(?:WS|AS)\s*:", line_clean, flags=re.I):
        return False

    return False


def extract_prices(text: str, aggressive: bool = False) -> tuple[float | None, float | None]:
    """
    Extract wholesale and retail prices.
    Most product lines are like:
        IR4565 | $70.00 WHOLESALE | $140.00 RETAIL | NEW
    Fallback keeps the first two dollar values as wholesale/retail.
    
    Args:
        text: Text to extract prices from
        aggressive: If True, will assign first smaller price to wholesale and larger to retail
                   even if labels are missing or unclear
    """
    text_clean = clean_text(text)

    wholesale = None
    retail = None
    debug_mode = False  # Set to True to see extraction details

    # Explicit label near value: $70.00 WHOLESALE / $140.00 RETAIL or $70.00Wholesale
    for match in re.finditer(
        r"\$\s*(\d+(?:\.\d{2})?)\s*(WHOLESALE|WS|WHSE|RETAIL)\b(?!\s*:)",
        text_clean,
        flags=re.I,
    ):
        value = float(match.group(1))
        label = match.group(2).upper()
        if label in {"WHOLESALE", "WS", "WHSE"} and wholesale is None:
            wholesale = value
            if debug_mode: print(f"  Found WHOLESALE (after $): {value}")
        elif label == "RETAIL" and retail is None:
            retail = value
            if debug_mode: print(f"  Found RETAIL (after $): {value}")

    # Explicit label before value: RETAIL $26.00 or Retail: $26.00
    for match in re.finditer(
        r"\b(WHOLESALE|WS|WHSE|RETAIL)\b\s*:?\s*\$\s*(\d+(?:\.\d{2})?)",
        text_clean,
        flags=re.I,
    ):
        label = match.group(1).upper()
        value = float(match.group(2))
        if label in {"WHOLESALE", "WS", "WHSE"} and wholesale is None:
            wholesale = value
            if debug_mode: print(f"  Found WHOLESALE (before $): {value}")
        elif label == "RETAIL" and retail is None:
            retail = value
            if debug_mode: print(f"  Found RETAIL (before $): {value}")

    # Same labeled formats without a dollar sign: 70.00 WHOLESALE / Retail: 140.00
    for match in re.finditer(
        r"(?<![A-Z0-9$])(\d+\.\d{2})\s*(WHOLESALE|WS|WHSE|RETAIL)\b(?!\s*:)",
        text_clean,
        flags=re.I,
    ):
        value = float(match.group(1))
        label = match.group(2).upper()
        if label in {"WHOLESALE", "WS", "WHSE"} and wholesale is None:
            wholesale = value
            if debug_mode: print(f"  Found WHOLESALE (no $ after value): {value}")
        elif label == "RETAIL" and retail is None:
            retail = value
            if debug_mode: print(f"  Found RETAIL (no $ after value): {value}")

    for match in re.finditer(
        r"\b(WHOLESALE|WS|WHSE|RETAIL)\b\s*:?\s*(?!\$)(\d+\.\d{2})",
        text_clean,
        flags=re.I,
    ):
        label = match.group(1).upper()
        value = float(match.group(2))
        if label in {"WHOLESALE", "WS", "WHSE"} and wholesale is None:
            wholesale = value
            if debug_mode: print(f"  Found WHOLESALE (no $ before value): {value}")
        elif label == "RETAIL" and retail is None:
            retail = value
            if debug_mode: print(f"  Found RETAIL (no $ before value): {value}")

    # Fallback: all explicit dollar amounts in order.
    values: list[float] = []
    for match in re.finditer(r"\$\s*(\d+(?:\.\d{2})?)", text_clean):
        value = float(match.group(1))
        if value not in values:
            values.append(value)
    
    if debug_mode and values: print(f"  All dollar values found: {values}")

    # Backup: values like 20.00 WS without $
    if len(values) < 2:
        for match in re.finditer(
            r"(?<![A-Z0-9])(\d+\.\d{2})(?=\s*(?:WS|WHSE|WHOLESALE|RETAIL|\|))",
            text_clean,
            flags=re.I,
        ):
            value = float(match.group(1))
            if value not in values:
                values.append(value)

    # First pass fallback for completely unlabeled prices only.
    # Do not copy a single explicitly labeled retail price into wholesale.
    if wholesale is None and retail is None:
        if len(values) >= 1:
            wholesale = values[0]
            if debug_mode: print(f"  Fallback WHOLESALE: {wholesale}")
        if len(values) >= 2:
            retail = values[1]
            if debug_mode: print(f"  Fallback RETAIL: {retail}")
    elif retail is None and len(values) >= 2:
        for value in values:
            if value != wholesale:
                retail = value
                if debug_mode: print(f"  Fallback RETAIL: {retail}")
                break
    
    if debug_mode:
        print(f"  After initial extraction - WS: {wholesale}, RT: {retail}")
    
    # Aggressive mode: intelligently assign prices by value when labels are unclear
    if aggressive and len(values) >= 2:
        unique_values = sorted(set(values))
        if debug_mode: print(f"  Aggressive mode: unique sorted values: {unique_values}")
        
        # If we found two different prices, assign smaller to wholesale, larger to retail
        if len(unique_values) >= 2:
            # Make sure wholesale is the smaller price
            if wholesale is None or retail is None:
                # At least one is missing, fill both intelligently
                wholesale = unique_values[0]  # Smaller
                retail = unique_values[1]     # Larger
                if debug_mode: print(f"  Aggressive: filled missing - WS: {wholesale}, RT: {retail}")
            elif wholesale > retail:
                # Prices are backwards, swap them
                wholesale, retail = retail, wholesale
                if debug_mode: print(f"  Aggressive: swapped backwards - WS: {wholesale}, RT: {retail}")

    if debug_mode:
        print(f"  Final result - WS: {wholesale}, RT: {retail}\n")

    return wholesale, retail


def is_likely_product_name(line: str) -> bool:
    line = clean_text(line)

    if not line:
        return False

    if STYLE_RE.search(line):
        return False

    if STOP_RE.search(line):
        return False

    if BAD_PRODUCT_NAME_RE.search(line):
        return False

    if FEATURE_PRODUCT_NAME_RE.search(line):
        return False

    if HEADER_PRODUCT_NAME_RE.search(line):
        return False

    # Color-code-only line: 010 or 010 100 274
    if re.fullmatch(r"[0-9]{3}(\s+[0-9]{3})*", line):
        return False

    # Size ranges like 3.5-13,14,15,16
    if re.fullmatch(r"\d+(?:\.\d+)?-\d+.*", line):
        return False

    if looks_like_price_context(line):
        return False

    # Reject obvious paragraph text.
    if len(line.split()) > 12 and not PRODUCT_TITLE_RE.search(line):
        return False

    return product_title_score(line) >= 20


def product_title_score(line: str) -> int:
    line_clean = clean_text(line)
    line_upper = line_clean.upper()
    tokens = set(re.findall(r"[A-Z0-9%'/.-]+", line_upper))

    score = 0

    # Strong Nike product title pattern
    if PRODUCT_TITLE_RE.search(line_clean):
        score += 100

    # Product category words
    if tokens & PRODUCT_KEYWORDS:
        score += 25

    # Many product names are uppercase
    if line_clean == line_upper and len(tokens) >= 3:
        score += 10

    # Product names are usually compact, not long descriptions.
    if 2 <= len(tokens) <= 10:
        score += 5

    # Penalize weak words that are often features/materials.
    if BAD_PRODUCT_NAME_RE.search(line_clean) or FEATURE_PRODUCT_NAME_RE.search(line_clean):
        score -= 100

    return score


def is_category_heading(line: str) -> bool:
    line_clean = clean_text(line)
    return any(pattern.fullmatch(line_clean) for _, pattern in CATEGORY_PATTERNS)


def accessory_slide_product_name(lines: list[str]) -> str | None:
    """
    Bag/accessory slides usually have one large product title near the top.
    Use that title to ignore hidden/leftover copied captions from another product.
    """
    skip_exact = {
        "STANDBAGS",
        "STAND BAGS",
        "BAGS",
        "ACCESSORIES",
        "KEEP IT TIGHT",
        "KEEP IT TIGHT.",
    }

    for line in lines[:25]:
        line_clean = clean_text(line)
        line_upper = line_clean.upper()

        if not line_clean:
            continue
        if line_upper in skip_exact or "NIKE GLOBAL" in line_upper:
            continue
        if is_category_heading(line_clean):
            continue
        if STYLE_RE.search(line_clean) or ACCESSORY_STYLE_COLOR_RE.search(line_clean):
            continue
        if looks_like_price_context(line_clean):
            continue
        if PRODUCT_TITLE_RE.search(line_clean) and is_likely_product_name(line_clean):
            return line_clean

    return None


def same_product_name(left: str | None, right: str | None) -> bool:
    if not left or not right:
        return False
    return clean_text(left).upper() == clean_text(right).upper()


def item_right(item: dict) -> int:
    return int(item.get("left", 0)) + int(item.get("width", 0))


def item_bottom(item: dict) -> int:
    return int(item.get("top", 0)) + int(item.get("height", 0))


def horizontal_overlap_ratio(a: dict, b: dict) -> float:
    overlap = max(0, min(item_right(a), item_right(b)) - max(int(a.get("left", 0)), int(b.get("left", 0))))
    base = max(1, min(int(a.get("width", 0)) or 1, int(b.get("width", 0)) or 1))
    return overlap / base


def is_grid_product_name_candidate(item: dict) -> bool:
    text = clean_text(item.get("text", ""))
    if not is_likely_product_name(text):
        return False
    if is_category_heading(text):
        return False
    if ACCESSORY_STYLE_COLOR_RE.search(text) or STYLE_COLOR_RE.search(text):
        return False
    return product_title_score(text) >= 10


def nearest_product_name_by_position(code_item: dict, items: list[dict]) -> tuple[str | None, float | None, bool]:
    candidates: list[dict] = []
    code_top = int(code_item.get("top", 0))
    code_center_x = int(code_item.get("left", 0)) + max(int(code_item.get("width", 0)), 1) / 2

    for item in items:
        if not is_grid_product_name_candidate(item):
            continue

        same_box_before = (
            item.get("shape_index") == code_item.get("shape_index")
            and int(item.get("line_index", 0)) < int(code_item.get("line_index", 0))
        )
        vertical_gap = code_top - item_bottom(item)
        same_column = horizontal_overlap_ratio(code_item, item) >= 0.25
        center_gap = abs((int(item.get("left", 0)) + max(int(item.get("width", 0)), 1) / 2) - code_center_x)

        if not same_box_before and (vertical_gap < -40_000 or vertical_gap > 1_200_000):
            continue
        if not same_box_before and not same_column and center_gap > max(int(code_item.get("width", 0)), int(item.get("width", 0)), 1):
            continue

        score = (
            product_title_score(item["text"])
            + (75 if same_box_before else 0)
            + (30 * horizontal_overlap_ratio(code_item, item))
            - (abs(vertical_gap) / 100_000)
            - (center_gap / 200_000)
        )
        candidates.append({"text": item["text"], "score": score, "same_box": same_box_before})

    if not candidates:
        return None, None, False

    best = max(candidates, key=lambda x: x["score"])
    return best["text"], round(float(best["score"]), 3), bool(best["same_box"])


def nearest_price_context_by_position(code_item: dict, items: list[dict]) -> tuple[str, float | None, float | None, str]:
    code_text = clean_text(code_item.get("text", ""))
    code_ws, code_retail = extract_prices(code_text, aggressive=True)
    if code_ws is not None and code_retail is not None:
        return code_text, code_ws, code_retail, "high"

    price_candidates: list[dict] = []
    code_top = int(code_item.get("top", 0))
    code_center_x = int(code_item.get("left", 0)) + max(int(code_item.get("width", 0)), 1) / 2

    for item in items:
        text = clean_text(item.get("text", ""))
        if not text or item is code_item or not looks_like_price_context(text):
            continue

        wholesale, retail = extract_prices(text, aggressive=True)
        if wholesale is None and retail is None:
            continue

        same_box = item.get("shape_index") == code_item.get("shape_index")
        item_center_x = int(item.get("left", 0)) + max(int(item.get("width", 0)), 1) / 2
        center_gap = abs(item_center_x - code_center_x)
        vertical_gap = int(item.get("top", 0)) - code_top
        overlap = horizontal_overlap_ratio(code_item, item)
        same_card = same_box or (
            -80_000 <= vertical_gap <= 650_000
            and overlap >= 0.45
            and center_gap <= max(int(code_item.get("width", 0)), int(item.get("width", 0)), 1) * 0.55
        )

        if same_card:
            score = (100 if same_box else 0) + (overlap * 50) - (abs(vertical_gap) / 100_000) - (center_gap / 200_000)
            price_candidates.append(
                {
                    "text": text,
                    "wholesale": wholesale,
                    "retail": retail,
                    "score": score,
                    "same_box": same_box,
                }
            )

    if not price_candidates:
        return code_text, None, None, "low"

    best = max(price_candidates, key=lambda x: x["score"])
    confidence = "high" if best["same_box"] and best["wholesale"] is not None and best["retail"] is not None else "medium"
    return clean_text(f"{code_text} {best['text']}"), best["wholesale"], best["retail"], confidence


def previous_product_name(lines: list[str], idx: int) -> str | None:
    candidates: list[dict] = []

    # Search farther back because some slides have:
    # title -> material/key features -> style code
    for j in range(idx - 1, max(-1, idx - 60), -1):
        line = lines[j]

        # Stop if we hit a previous product block.
        if STYLE_RE.search(line) and looks_like_price_context(line):
            break

        if is_likely_product_name(line):
            candidates.append(
                {
                    "line": line,
                    "score": product_title_score(line),
                    "distance": abs(idx - j),
                }
            )

    if not candidates:
        return None

    candidates = sorted(
        candidates,
        key=lambda x: (x["score"], -x["distance"]),
        reverse=True,
    )

    return candidates[0]["line"]


def next_product_name(lines: list[str], idx: int) -> str | None:
    candidates: list[dict] = []

    for j in range(idx + 1, min(len(lines), idx + 8)):
        line = lines[j]

        if STYLE_RE.search(line) and looks_like_price_context(line):
            break

        if is_likely_product_name(line):
            candidates.append(
                {
                    "line": line,
                    "score": product_title_score(line),
                    "distance": abs(idx - j),
                }
            )

    if not candidates:
        return None

    candidates = sorted(
        candidates,
        key=lambda x: (x["score"], -x["distance"]),
        reverse=True,
    )

    return candidates[0]["line"]


def build_product_name(line: str, lines: list[str], idx: int, first_style: str) -> str | None:
    pos = line.find(first_style)

    # Case where product name and style code are on the same line:
    # "M NK DF PAR POLO SS SOLID IB0233 | $45..."
    prefix = clean_text(line[:pos].replace("Style", "")) if pos > 0 else ""

    if prefix and is_likely_product_name(prefix):
        return prefix

    return previous_product_name(lines, idx) or next_product_name(lines, idx) or None


def choose_recovery_name(rows: pd.DataFrame) -> dict[str, str]:
    if rows.empty or "style_code" not in rows.columns or "product_name" not in rows.columns:
        return {}

    candidates: dict[str, tuple[int, str]] = {}
    for _, row in rows.iterrows():
        style = clean_text(row.get("style_code", "")).upper()
        name = clean_text(row.get("product_name", ""))
        if not style or not is_likely_product_name(name):
            continue

        score = product_title_score(name) + len(name.split())
        if style not in candidates or score > candidates[style][0]:
            candidates[style] = (score, name)

    return {style: value[1] for style, value in candidates.items()}


def choose_preferred_name_by_key(rows: pd.DataFrame) -> dict[str, str]:
    if rows.empty or "style_color_key" not in rows.columns or "product_name" not in rows.columns:
        return {}

    candidates: dict[str, tuple[int, str]] = {}
    for _, row in rows.iterrows():
        key = clean_text(row.get("style_color_key", "")).upper()
        name = clean_text(row.get("product_name", ""))
        if not key or not is_likely_product_name(name):
            continue

        source = clean_text(row.get("extraction_source", ""))
        source_bonus = 50 if source == "detail_product_parser" else 0
        score = source_bonus + product_title_score(name) + len(name)
        if key not in candidates or score > candidates[key][0]:
            candidates[key] = (score, name)

    return {key: value[1] for key, value in candidates.items()}


def item_is_on_slide(item: dict, slide_width: int, slide_height: int) -> bool:
    left = int(item.get("left", 0))
    top = int(item.get("top", 0))
    width = int(item.get("width", 0))
    height = int(item.get("height", 0))
    return not (
        left + width < 0
        or top + height < 0
        or left > slide_width
        or top > slide_height
    )


# -----------------------------
# Nike accessories / bags format
# -----------------------------

def split_accessory_style_color(full_code: str) -> tuple[str, str]:
    """
    Converts Nike accessory/bag codes like N.100.3478.091 into:
        style_code = N.100.3478
        color_code = 091
    """
    code = clean_text(full_code).upper()
    parts = code.split(".")
    if len(parts) == 4:
        return ".".join(parts[:3]), parts[3].zfill(3)
    return code, ""


def has_accessory_style_color(text: str) -> bool:
    return bool(ACCESSORY_STYLE_COLOR_RE.search(clean_text(text)))


def get_text_box_lines(items: list[dict], shape_index: int) -> list[str]:
    shape_items = [item for item in items if item["shape_index"] == shape_index]
    shape_items = sorted(shape_items, key=lambda x: x.get("line_index", 0))
    return [item["text"] for item in shape_items]


def accessory_product_name_from_same_box(items: list[dict], code_item: dict) -> str | None:
    shape_items = [
        item
        for item in items
        if item["shape_index"] == code_item["shape_index"]
        and int(item.get("line_index", 0)) < int(code_item.get("line_index", 0))
    ]
    shape_items = sorted(shape_items, key=lambda x: x.get("line_index", 0), reverse=True)

    for item in shape_items:
        line = clean_text(item["text"])
        if is_likely_product_name(line):
            return line

    return None


def get_accessory_price_context(items: list[dict], code_item: dict, lines: list[str], line_idx: int) -> str:
    """
    Accessory slides often place each product caption in one text box:
        product name
        N.101.4761.016
        Wholesale: $17.50
        Retail: $35.00

    Use that text box first so prices from neighboring captions are not mixed in.
    """
    same_box_lines = get_text_box_lines(items, code_item["shape_index"])
    same_box_text = clean_text(" ".join(same_box_lines))
    wholesale, retail = extract_prices(same_box_text, aggressive=True)

    if wholesale is not None and retail is not None:
        return same_box_text

    left = int(code_item.get("left", 0))
    top = int(code_item.get("top", 0))
    width = max(int(code_item.get("width", 0)), 1)
    height = max(int(code_item.get("height", 0)), 1)
    center_x = left + width / 2

    nearby: list[dict] = []
    for item in items:
        item_left = int(item.get("left", 0))
        item_top = int(item.get("top", 0))
        item_width = max(int(item.get("width", 0)), 1)
        item_center_x = item_left + item_width / 2

        same_column = abs(item_center_x - center_x) <= max(width, item_width) * 0.65
        close_below = top - height <= item_top <= top + height * 3.5

        if same_column and close_below:
            nearby.append(item)

    nearby_text = clean_text(" ".join(item["text"] for item in sorted(nearby, key=lambda x: (x["top"], x["left"], x.get("line_index", 0)))))
    wholesale, retail = extract_prices(nearby_text, aggressive=True)
    if wholesale is not None and retail is not None:
        return nearby_text

    return same_box_text or nearby_text or clean_text(" ".join(lines[line_idx : min(line_idx + 10, len(lines))]))


def build_style_color_key(style: str, color: str | None) -> str:
    style = clean_text(style)
    color = clean_text(color) if color is not None else ""

    if not color:
        return style

    # Accessories already use dot notation in the original Nike code.
    if style.upper().startswith("N."):
        return f"{style}.{color}"

    # Apparel uses the regular style-color convention.
    return f"{style}-{color}"


def build_image_name(slide_number: int, style_color_key: str) -> str:
    safe_key = re.sub(r"[^A-Za-z0-9._-]+", "_", clean_text(style_color_key))
    return f"slide_{int(slide_number):03d}_{safe_key}.png"


def iter_picture_shapes(shapes):
    """Yield picture-like shapes recursively, including pictures inside groups."""
    for shape in shapes:
        if hasattr(shape, "image"):
            yield shape
        if hasattr(shape, "shapes"):
            yield from iter_picture_shapes(shape.shapes)


def picture_to_png_bytes(blob: bytes) -> bytes | None:
    try:
        with Image.open(io.BytesIO(blob)) as image:
            if image.mode not in {"RGB", "RGBA"}:
                image = image.convert("RGBA")
            output = io.BytesIO()
            image.save(output, format="PNG")
            return output.getvalue()
    except (UnidentifiedImageError, OSError):
        return None


def get_slide_product_pictures(slide) -> list[dict]:
    pictures: list[dict] = []

    for shape in iter_picture_shapes(slide.shapes):
        width = int(getattr(shape, "width", 0))
        height = int(getattr(shape, "height", 0))

        # Skip tiny logos/icons. Product images are materially larger than these.
        if width < 300_000 or height < 300_000:
            continue

        pictures.append(
            {
                "blob": shape.image.blob,
                "left": int(getattr(shape, "left", 0)),
                "top": int(getattr(shape, "top", 0)),
                "width": width,
                "height": height,
            }
        )

    return sorted(pictures, key=lambda x: (x["top"], x["left"]))


def build_images_zip_from_pptx(pptx_path: str | Path, style_colors: pd.DataFrame) -> tuple[bytes, int]:
    """
    Export product pictures into a ZIP using the image_name column.
    Matching is by slide and visual order, so it is a practical database seed rather than
    a pixel-perfect image/product matcher.
    """
    if style_colors.empty or "image_name" not in style_colors.columns:
        return b"", 0

    prs = Presentation(str(pptx_path))
    zip_buffer = io.BytesIO()
    written_names: set[str] = set()
    image_count = 0

    with zipfile.ZipFile(zip_buffer, "w", compression=zipfile.ZIP_DEFLATED) as zip_file:
        for slide_number, slide in enumerate(prs.slides, start=1):
            slide_rows = style_colors[style_colors["slide_number"] == slide_number].copy()
            if slide_rows.empty:
                continue

            pictures = get_slide_product_pictures(slide)
            if not pictures:
                continue

            slide_rows = slide_rows.drop_duplicates("image_name").reset_index(drop=True)

            for row_idx, row in slide_rows.iterrows():
                picture_idx = min(
                    int(row_idx * len(pictures) / max(len(slide_rows), 1)),
                    len(pictures) - 1,
                )
                png_bytes = picture_to_png_bytes(pictures[picture_idx]["blob"])
                if png_bytes is None:
                    continue

                image_name = str(row["image_name"])
                if image_name in written_names:
                    continue

                zip_file.writestr(image_name, png_bytes)
                written_names.add(image_name)
                image_count += 1

    zip_buffer.seek(0)
    return zip_buffer.getvalue(), image_count


# -----------------------------
# Color extraction
# -----------------------------

def is_color_code_line(line: str, slide_number: int | None = None) -> bool:
    line = clean_text(line)

    if not line:
        return False

    if slide_number is not None and line == str(slide_number):
        return False

    if STYLE_RE.search(line):
        return False

    if looks_like_price_context(line):
        return False

    if STOP_RE.search(line):
        return False

    # Ignore size ranges
    if re.fullmatch(r"\d+(?:\.\d+)?-\d+.*", line):
        return False

    # Ignore years
    if re.search(r"\b20\d{2}\b", line):
        return False

    # Exact color code: 010
    if re.fullmatch(r"\d{3}", line):
        return True

    # Multiple exact color codes: 010 100 274
    if re.fullmatch(r"(\d{3}\s*)+", line):
        return True

    # Color code with description: 010 - BLACK, 010 BLACK
    if re.match(r"^\d{3}\s*(?:[-/]\s*)?[A-Za-z]", line):
        return True

    return False


def extract_color_codes_from_text(line: str, slide_number: int | None = None) -> list[str]:
    line = clean_text(line)

    if not is_color_code_line(line, slide_number=slide_number):
        return []

    # Exact code / multiple exact codes
    if re.fullmatch(r"(\d{3}\s*)+", line):
        return [code.zfill(3) for code in re.findall(r"\d{3}", line)]

    # Code + description
    match = re.match(r"^(\d{3})\s*(?:[-/]\s*)?[A-Za-z]", line)
    if match:
        return [match.group(1).zfill(3)]

    return []


def extract_color_codes_from_lines(lines: list[str], start_idx: int) -> list[str]:
    """
    Extract color codes near the product block.
    This is useful, but not always enough because PowerPoint shape order can be strange.
    """
    color_codes: list[str] = []

    for line in lines[start_idx : start_idx + 30]:
        if STOP_RE.search(line):
            break

        if STYLE_RE.search(line) and looks_like_price_context(line):
            break

        color_codes.extend(extract_color_codes_from_text(line))

    return unique_preserve_order(color_codes)


def extract_color_codes_from_slide_shapes(
    slide,
    slide_number: int,
    slide_width: int | None = None,
    slide_height: int | None = None,
) -> list[str]:
    """
    Extract color codes from all text boxes on the slide.
    This catches separate labels under product images, like 010, 274, 393, 419.
    """
    color_codes: list[str] = []
    items = get_slide_text_items(slide, slide_width, slide_height)

    for item in items:
        color_codes.extend(extract_color_codes_from_text(item["text"], slide_number=slide_number))

    return unique_preserve_order(color_codes)


# -----------------------------
# Section extraction
# -----------------------------

def extract_section(lines: list[str], start_idx: int, label: str, stop_labels: Iterable[str]) -> str:
    results: list[str] = []
    active = False
    label_norm = label.lower().rstrip(":")
    stop_norm = tuple(s.lower().rstrip(":") for s in stop_labels)

    for line in lines[start_idx : start_idx + 60]:
        line_norm = line.lower().rstrip(":")

        if active:
            if any(line_norm.startswith(s) for s in stop_norm):
                break
            if STYLE_RE.search(line) and looks_like_price_context(line):
                break
            results.append(line)

        elif line_norm.startswith(label_norm):
            after_label = line[len(label) :].strip()
            if after_label:
                results.append(after_label)
            active = True

    return " | ".join(results)


def parse_grid_assortment_slide(
    slide_number: int,
    category: str | None,
    text_items: list[dict],
) -> list[dict]:
    """
    Parse suggested-look/assortment cards by spatial proximity instead of slide text order.
    """
    rows: list[dict] = []

    for item in text_items:
        expanded_codes = expand_standard_style_color_tokens(item["text"])
        if not expanded_codes:
            continue

        product_name, name_score, name_same_box = nearest_product_name_by_position(item, text_items)
        context, wholesale, retail, price_confidence = nearest_price_context_by_position(item, text_items)
        status_match = STATUS_RE.search(context)
        status = status_match.group(1).upper() if status_match else None

        confidence = (
            "high"
            if product_name and name_same_box and wholesale is not None and retail is not None and price_confidence == "high"
            else "medium"
        )
        if product_name is None or wholesale is None or retail is None:
            confidence = "low"
        elif price_confidence != "high":
            confidence = "medium"

        for code in expanded_codes:
            rows.append(
                {
                    "slide_number": slide_number,
                    "category": category,
                    "product_name": product_name,
                    "style_codes": code["style_code"],
                    "wholesale_price": wholesale,
                    "retail_price": retail,
                    "status": status,
                    "color_codes": code["color_code"],
                    "color_count": 1,
                    "material": "",
                    "key_features": "",
                    "raw_product_line": context,
                    "original_raw_code": code["original_raw_code"],
                    "slash_color_expanded": code["slash_color_expanded"],
                    "extraction_source": "grid_position_parser",
                    "confidence": confidence,
                    "name_match_score": name_score,
                }
            )

    return rows


# -----------------------------
# Main extraction
# -----------------------------

def extract_products_from_pptx(
    pptx_path: str | Path,
) -> tuple[pd.DataFrame, pd.DataFrame, pd.DataFrame, pd.DataFrame]:
    prs = Presentation(str(pptx_path))

    raw_rows: list[dict] = []
    product_rows: list[dict] = []
    used_line_indices: set[tuple[int, int]] = set()
    current_category: str | None = None

    for slide_number, slide in enumerate(prs.slides, start=1):
        slide_width = int(prs.slide_width)
        slide_height = int(prs.slide_height)
        text_items = get_slide_text_items(slide, slide_width, slide_height)
        all_text_items = get_slide_text_items(slide)
        lines = [item["text"] for item in text_items]
        detected_category = detect_slide_category(lines)
        if detected_category:
            current_category = detected_category
        primary_accessory_product_name = (
            accessory_slide_product_name(lines)
            if current_category in {"Bags", "Gloves", "Headcovers", "Accessories"}
            else None
        )
        raw_rows.append(
            {
                "slide_number": slide_number,
                "category": current_category,
                "slide_text": "\n".join(lines),
            }
        )

        # Find all product-start lines on the slide.
        product_start_indices: list[int] = []
        for idx, slide_line in enumerate(lines):
            lookahead = " ".join(lines[idx + 1 : idx + 3])
            if is_new_product_start(slide_line, lookahead):
                product_start_indices.append(idx)

        slide_shape_colors = extract_color_codes_from_slide_shapes(
            slide,
            slide_number,
            slide_width,
            slide_height,
        )

        grid_rows = parse_grid_assortment_slide(slide_number, current_category, text_items)
        grid_style_color_keys = {
            build_style_color_key(row["style_codes"], row["color_codes"])
            for row in grid_rows
            if row.get("style_codes") and row.get("color_codes")
        }
        product_rows.extend(grid_rows)

        # Special handling for Nike accessories/bags slides.
        # These use codes like N.100.3478.091 and prices on the following lines.
        accessory_line_items = {
            (item["shape_index"], item.get("line_index", 0)): idx
            for idx, item in enumerate(all_text_items)
        }

        for item in all_text_items:
            line = item["text"]
            i = accessory_line_items.get((item["shape_index"], item.get("line_index", 0)), 0)
            accessory_codes = ACCESSORY_STYLE_COLOR_RE.findall(line)
            if not accessory_codes:
                continue

            same_box_product_name = accessory_product_name_from_same_box(all_text_items, item)
            on_slide = item_is_on_slide(item, slide_width, slide_height)
            if (
                not on_slide
                and primary_accessory_product_name
                and not same_product_name(same_box_product_name, primary_accessory_product_name)
            ):
                continue

            combined = get_accessory_price_context(all_text_items, item, lines, i)
            # Use aggressive mode for bags to find both prices even if labels are unclear
            wholesale, retail = extract_prices(combined, aggressive=True)
            status_match = STATUS_RE.search(combined)
            status = status_match.group(1).upper() if status_match else None
            spatial_product_name, name_score, _ = nearest_product_name_by_position(item, text_items)
            product_name = same_box_product_name or spatial_product_name or previous_product_name(lines, i) or next_product_name(lines, i)
            if not is_likely_product_name(product_name or "") and primary_accessory_product_name:
                product_name = primary_accessory_product_name

            material = extract_section(
                lines,
                i + 1,
                "Material Content",
                ["Features and Benefits", "Key Features:", "Overview:", "Product Details:"],
            )

            key_features = extract_section(
                lines,
                i + 1,
                "Features and Benefits",
                ["Material Content", "Material:", "Overview:", "Product Details:"],
            )

            for full_code in accessory_codes:
                style_code, color_code = split_accessory_style_color(full_code)

                product_rows.append(
                    {
                        "slide_number": slide_number,
                        "category": current_category,
                        "product_name": product_name,
                        "style_codes": style_code,
                        "wholesale_price": wholesale,
                        "retail_price": retail,
                        "status": status,
                        "color_codes": color_code,
                        "color_count": 1 if color_code else 0,
                        "material": material,
                        "key_features": key_features,
                        "raw_product_line": combined,
                        "original_raw_code": full_code.upper(),
                        "slash_color_expanded": False,
                        "extraction_source": "bag_accessory_parser",
                        "confidence": "high" if same_box_product_name and wholesale is not None and retail is not None else "medium",
                        "name_match_score": name_score,
                    }
                )

            used_line_indices.add((slide_number, i))

        for i, line in enumerate(lines):
            lookahead = " ".join(lines[i + 1 : i + 3])

            if (slide_number, i) in used_line_indices:
                continue

            detail_line = parse_detail_style_line(line)
            if detail_line is None:
                continue

            first_style = detail_line["style"]
            end_idx = i
            combined = line
            style_codes = [first_style]
            wholesale = detail_line["wholesale"]
            retail = detail_line["retail"]
            status = detail_line["status"]

            product_name = build_product_name(line, lines, i, first_style)
            if not is_likely_product_name(product_name or ""):
                spatial_product_name, _, _ = nearest_product_name_by_position(text_items[i], text_items)
                if spatial_product_name:
                    product_name = spatial_product_name

            # Original line-based extraction.
            line_colors = unique_preserve_order(
                extract_color_codes_from_lines(lines, end_idx + 1)
            )

            # If the slide has only one product, use all color labels found on the slide.
            # This catches colors stored as separate text boxes.
            if len(product_start_indices) <= 1:
                color_codes = unique_preserve_order(line_colors + slide_shape_colors)
            else:
                # Multi-product slides are harder.
                # Do NOT assign all slide colors to every product, because that creates false extra rows.
                # Keep only the safer line-based colors.
                color_codes = line_colors

            material = extract_section(
                lines,
                end_idx + 1,
                "Material:",
                ["Key Features:", "Overview:", "Product Details:"],
            )

            key_features = extract_section(
                lines,
                end_idx + 1,
                "Key Features:",
                ["Overview:", "Material:", "Product Details:"],
            )

            product_rows.append(
                {
                    "slide_number": slide_number,
                    "category": current_category,
                    "product_name": product_name,
                    "style_codes": ",".join(style_codes),
                    "wholesale_price": wholesale,
                    "retail_price": retail,
                    "status": status,
                    "color_codes": ",".join(color_codes),
                    "color_count": len(color_codes),
                    "material": material,
                    "key_features": key_features,
                    "raw_product_line": combined,
                    "original_raw_code": "",
                    "slash_color_expanded": False,
                    "extraction_source": "detail_product_parser",
                    "confidence": "high" if product_name and wholesale is not None and retail is not None else "medium",
                    "name_match_score": None,
                }
            )

            for j in range(i, end_idx + 1):
                used_line_indices.add((slide_number, j))

    raw_extract_df = pd.DataFrame(product_rows)
    raw_slide_text_df = pd.DataFrame(raw_rows)

    exploded_rows: list[dict] = []

    for _, row in raw_extract_df.iterrows():
        styles = [x for x in str(row["style_codes"]).split(",") if x]
        colors = [x for x in str(row["color_codes"]).split(",") if x]

        if not styles:
            continue

        if colors:
            for style in styles:
                for color in colors:
                    style_color_key = build_style_color_key(style, color)
                    exploded_rows.append(
                        {
                            "slide_number": row["slide_number"],
                            "category": row.get("category"),
                            "product_name": row["product_name"],
                            "style_code": style,
                            "color_code": color,
                            "style_color_key": style_color_key,
                            "image_name": build_image_name(row["slide_number"], style_color_key),
                            "wholesale_price": row["wholesale_price"],
                            "retail_price": row["retail_price"],
                            "status": row["status"],
                            "original_raw_code": row.get("original_raw_code", ""),
                            "slash_color_expanded": bool(row.get("slash_color_expanded", False)),
                            "extraction_source": row.get("extraction_source", "detail_product_parser"),
                            "confidence": row.get("confidence", "medium"),
                            "duplicate_flag": False,
                            "mismatch_flag": False,
                        }
                    )
        else:
            for style in styles:
                style_color_key = style
                exploded_rows.append(
                    {
                        "slide_number": row["slide_number"],
                        "category": row.get("category"),
                        "product_name": row["product_name"],
                        "style_code": style,
                        "color_code": None,
                        "style_color_key": style_color_key,
                        "image_name": build_image_name(row["slide_number"], style_color_key),
                        "wholesale_price": row["wholesale_price"],
                        "retail_price": row["retail_price"],
                        "status": row["status"],
                        "original_raw_code": row.get("original_raw_code", ""),
                        "slash_color_expanded": bool(row.get("slash_color_expanded", False)),
                        "extraction_source": row.get("extraction_source", "detail_product_parser"),
                        "confidence": row.get("confidence", "medium"),
                        "duplicate_flag": False,
                        "mismatch_flag": False,
                    }
                )

    raw_style_color_df = pd.DataFrame(exploded_rows)
    qa_rows: list[dict] = []

    if raw_style_color_df.empty:
        products_df = raw_extract_df
        style_color_df = raw_style_color_df
    else:
        valid_mask = raw_style_color_df["style_color_key"].astype(str).map(is_valid_style_color_key)
        invalid_rows = raw_style_color_df[~valid_mask].copy()
        for _, invalid in invalid_rows.iterrows():
            qa_rows.append(
                {
                    "slide_number": invalid.get("slide_number"),
                    "style_color_key": invalid.get("style_color_key"),
                    "issue_type": "invalid_style_color_key",
                    "message": (
                        "Rejected row because style_color_key does not match "
                        "AA0000-000 or N.000.0000.000 format."
                    ),
                    "extraction_source": invalid.get("extraction_source"),
                    "confidence": invalid.get("confidence"),
                }
            )

        raw_style_color_df = raw_style_color_df[valid_mask].reset_index(drop=True)

    if raw_style_color_df.empty:
        products_df = pd.DataFrame(columns=raw_extract_df.columns)
        style_color_df = raw_style_color_df
    else:
        recovery_names = choose_recovery_name(raw_style_color_df)
        for idx, row in raw_style_color_df.iterrows():
            current_name = clean_text(row.get("product_name", ""))
            style = clean_text(row.get("style_code", "")).upper()
            recovered_name = recovery_names.get(style)
            if recovered_name and not is_likely_product_name(current_name):
                raw_style_color_df.at[idx, "product_name"] = recovered_name
                raw_style_color_df.at[idx, "confidence"] = "medium"
                qa_rows.append(
                    {
                        "slide_number": row.get("slide_number"),
                        "style_color_key": row.get("style_color_key"),
                        "issue_type": "recovered_product_name",
                        "message": f"Recovered product_name from another valid row with style_code {style}: {recovered_name}",
                        "extraction_source": row.get("extraction_source"),
                        "confidence": "medium",
                    }
                )

        name_valid_mask = raw_style_color_df["product_name"].fillna("").astype(str).map(is_likely_product_name)
        questionable_name_rows = raw_style_color_df[~name_valid_mask].copy()
        for _, invalid in questionable_name_rows.iterrows():
            qa_rows.append(
                {
                    "slide_number": invalid.get("slide_number"),
                    "style_color_key": invalid.get("style_color_key"),
                    "issue_type": "invalid_or_missing_product_name",
                    "message": "Rejected row from clean output because product_name is missing or looks like feature/material/category text.",
                    "extraction_source": invalid.get("extraction_source"),
                    "confidence": invalid.get("confidence"),
                }
            )

        raw_style_color_df = raw_style_color_df[name_valid_mask].reset_index(drop=True)

    if raw_style_color_df.empty:
        products_df = pd.DataFrame(columns=raw_extract_df.columns)
        style_color_df = raw_style_color_df
    else:
        preferred_names = choose_preferred_name_by_key(raw_style_color_df)
        for idx, row in raw_style_color_df.iterrows():
            key = clean_text(row.get("style_color_key", "")).upper()
            current_name = clean_text(row.get("product_name", ""))
            preferred_name = preferred_names.get(key)
            if preferred_name and preferred_name != current_name and len(preferred_name) > len(current_name):
                raw_style_color_df.at[idx, "product_name"] = preferred_name
                raw_style_color_df.at[idx, "confidence"] = "medium"
                qa_rows.append(
                    {
                        "slide_number": row.get("slide_number"),
                        "style_color_key": row.get("style_color_key"),
                        "issue_type": "preferred_detail_product_name",
                        "message": f"Preferred fuller product_name for style-color: {preferred_name}",
                        "extraction_source": row.get("extraction_source"),
                        "confidence": "medium",
                    }
                )

        raw_style_color_df["duplicate_flag"] = raw_style_color_df.duplicated(
            subset=["style_color_key", "product_name", "wholesale_price", "retail_price"],
            keep=False,
        )

        mismatch_keys = set()
        for key, group in raw_style_color_df.groupby("style_color_key", dropna=False):
            names = set(group["product_name"].fillna("").astype(str).str.strip())
            prices = set(
                zip(
                    group["wholesale_price"].fillna("").astype(str),
                    group["retail_price"].fillna("").astype(str),
                )
            )
            if len(names) > 1 or len(prices) > 1:
                mismatch_keys.add(key)
                qa_rows.append(
                    {
                        "slide_number": ",".join(map(str, sorted(group["slide_number"].astype(str).unique()))),
                        "style_color_key": key,
                        "issue_type": "duplicate_style_color_mismatch",
                        "message": "Same style-color appears with different product names or prices.",
                        "extraction_source": ",".join(sorted(group["extraction_source"].dropna().astype(str).unique())),
                        "confidence": "medium",
                    }
                )

        raw_style_color_df["mismatch_flag"] = raw_style_color_df["style_color_key"].isin(mismatch_keys)

        clean_subset = ["style_color_key", "product_name", "wholesale_price", "retail_price"]
        style_color_df = raw_style_color_df.drop_duplicates(subset=clean_subset, keep="first").reset_index(drop=True)
        products_df = (
            style_color_df.groupby(
                [
                    "slide_number",
                    "category",
                    "product_name",
                    "style_code",
                    "wholesale_price",
                    "retail_price",
                    "status",
                    "extraction_source",
                    "confidence",
                ],
                dropna=False,
                as_index=False,
            )
            .agg(
                color_codes=("color_code", lambda x: ",".join(unique_preserve_order([str(v) for v in x if pd.notna(v)]))),
                color_count=("color_code", lambda x: int(pd.Series(x).dropna().nunique())),
                raw_style_color_count=("style_color_key", "size"),
            )
            .rename(columns={"style_code": "style_codes"})
        )

    raw_style_iter = raw_style_color_df.iterrows() if not raw_style_color_df.empty else []
    for _, row in raw_style_iter:
        if not str(row.get("product_name", "") or "").strip():
            qa_rows.append(
                {
                    "slide_number": row.get("slide_number"),
                    "style_color_key": row.get("style_color_key"),
                    "issue_type": "missing_product_name",
                    "message": "Style-color code detected but no product name was confidently matched.",
                    "extraction_source": row.get("extraction_source"),
                    "confidence": row.get("confidence"),
                }
            )
        if pd.isna(row.get("wholesale_price")) or pd.isna(row.get("retail_price")):
            qa_rows.append(
                {
                    "slide_number": row.get("slide_number"),
                    "style_color_key": row.get("style_color_key"),
                    "issue_type": "missing_price",
                    "message": "Product/style-color detected but wholesale or retail price is missing.",
                    "extraction_source": row.get("extraction_source"),
                    "confidence": row.get("confidence"),
                }
            )
        if row.get("slash_color_expanded"):
            qa_rows.append(
                {
                    "slide_number": row.get("slide_number"),
                    "style_color_key": row.get("style_color_key"),
                    "issue_type": "slash_color_expanded",
                    "message": f"Expanded slash color from {row.get('original_raw_code')}.",
                    "extraction_source": row.get("extraction_source"),
                    "confidence": row.get("confidence"),
                }
            )
        if str(row.get("confidence", "")).lower() == "low":
            qa_rows.append(
                {
                    "slide_number": row.get("slide_number"),
                    "style_color_key": row.get("style_color_key"),
                    "issue_type": "low_confidence",
                    "message": "Spatial parser produced a low-confidence row.",
                    "extraction_source": row.get("extraction_source"),
                    "confidence": row.get("confidence"),
                }
            )

    extracted_keys = set(raw_style_color_df["style_color_key"]) if not raw_style_color_df.empty else set()

    for _, row in raw_slide_text_df.iterrows():
        qa_rows.extend(
            {
                "slide_number": row["slide_number"],
                "style_color_key": "",
                "issue_type": "unsupported_or_ambiguous_line",
                "message": line,
                "extraction_source": "raw_slide_text",
                "confidence": "low",
            }
            for line in str(row.get("slide_text", "")).splitlines()
            if (STYLE_COLOR_RE.search(line) or ACCESSORY_STYLE_COLOR_RE.search(line))
            and not any(code["style_color_key"] in extracted_keys for code in expand_standard_style_color_tokens(line))
            and not any(code.upper() in extracted_keys for code in ACCESSORY_STYLE_COLOR_RE.findall(line))
        )

    qa_issues_df = pd.DataFrame(qa_rows)
    raw_extract_df = raw_style_color_df if not raw_style_color_df.empty else raw_extract_df

    return products_df, style_color_df, raw_extract_df, qa_issues_df


# -----------------------------
# Optional validation report
# -----------------------------

def write_validation_report(style_colors: pd.DataFrame, expected_csv: Path, output_folder: Path) -> None:
    """
    Compare the extracted style-color output to a validation CSV.

    The validation CSV should have at least:
        slide_number, product_name, style_code, color_code, style_color_key
    """
    if not expected_csv.exists():
        print("No validation CSV found. Skipping validation report.")
        return

    expected = pd.read_csv(expected_csv)
    actual = style_colors.copy()

    for df in (expected, actual):
        df["slide_number"] = pd.to_numeric(df["slide_number"], errors="coerce").astype("Int64")
        df["style_color_key"] = df["style_color_key"].astype(str).str.strip()

    compare = expected.merge(
        actual,
        on=["slide_number", "style_color_key"],
        how="outer",
        suffixes=("_expected", "_actual"),
        indicator=True,
    )

    missing = compare[compare["_merge"] == "left_only"].copy()
    extra = compare[compare["_merge"] == "right_only"].copy()

    product_name_mismatch = compare[
        (compare["_merge"] == "both")
        & (
            compare["product_name_expected"].fillna("").astype(str).str.strip()
            != compare["product_name_actual"].fillna("").astype(str).str.strip()
        )
    ].copy()

    price_mismatch = compare[
        (compare["_merge"] == "both")
        & (
            (compare["wholesale_price_expected"].fillna(-1) != compare["wholesale_price_actual"].fillna(-1))
            | (compare["retail_price_expected"].fillna(-1) != compare["retail_price_actual"].fillna(-1))
        )
    ].copy()

    status_mismatch = compare[
        (compare["_merge"] == "both")
        & (
            compare["status_expected"].fillna("").astype(str).str.strip().str.upper()
            != compare["status_actual"].fillna("").astype(str).str.strip().str.upper()
        )
    ].copy()

    missing.to_csv(output_folder / "validation_missing_rows.csv", index=False)
    extra.to_csv(output_folder / "validation_extra_rows.csv", index=False)
    product_name_mismatch.to_csv(output_folder / "validation_product_name_mismatch.csv", index=False)
    price_mismatch.to_csv(output_folder / "validation_price_mismatch.csv", index=False)
    status_mismatch.to_csv(output_folder / "validation_status_mismatch.csv", index=False)

    print("\nValidation report:")
    print(f"Missing rows: {len(missing)}")
    print(f"Extra rows: {len(extra)}")
    print(f"Product-name mismatches: {len(product_name_mismatch)}")
    print(f"Price mismatches: {len(price_mismatch)}")
    print(f"Status mismatches: {len(status_mismatch)}")
