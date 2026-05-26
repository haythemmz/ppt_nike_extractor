from __future__ import annotations

import argparse
import csv
import json
import re
from datetime import date, datetime
from pathlib import Path
from typing import Any, Iterable

from pptx import Presentation


EMU_PER_INCH = 914400


def clean_text(value: str | None) -> str:
    if not value:
        return ""
    return re.sub(r"\s+", " ", value).strip()


def safe_filename(value: str, fallback: str = "untitled") -> str:
    value = clean_text(value) or fallback
    value = re.sub(r"[^A-Za-z0-9._-]+", "_", value)
    value = re.sub(r"_+", "_", value).strip("._-")
    return value[:80] or fallback


def emu_to_inches(value: int | None) -> float | None:
    if value is None:
        return None
    return round(int(value) / EMU_PER_INCH, 3)


def json_safe(value: Any) -> Any:
    if isinstance(value, (datetime, date)):
        return value.isoformat()
    return value


def iter_shapes(shapes) -> Iterable[Any]:
    for shape in shapes:
        yield shape
        if hasattr(shape, "shapes"):
            yield from iter_shapes(shape.shapes)


def iter_picture_shapes(shapes) -> Iterable[Any]:
    for shape in iter_shapes(shapes):
        if hasattr(shape, "image"):
            yield shape


def get_shape_text(shape) -> str:
    if not hasattr(shape, "text"):
        return ""
    return clean_text(shape.text)


def get_slide_text_items(slide) -> list[dict[str, Any]]:
    items: list[dict[str, Any]] = []
    for shape in iter_shapes(slide.shapes):
        text = get_shape_text(shape)
        if not text:
            continue
        items.append(
            {
                "text": text,
                "left": int(getattr(shape, "left", 0) or 0),
                "top": int(getattr(shape, "top", 0) or 0),
                "width": int(getattr(shape, "width", 0) or 0),
                "height": int(getattr(shape, "height", 0) or 0),
            }
        )
    return sorted(items, key=lambda item: (item["top"], item["left"]))


def get_slide_title(slide, text_items: list[dict[str, Any]]) -> str:
    title_shape = getattr(slide.shapes, "title", None)
    title = get_shape_text(title_shape) if title_shape is not None else ""
    if title:
        return title

    for item in text_items:
        text = item["text"]
        if len(text) > 2 and not text.isdigit():
            return text[:140]

    return ""


def get_position_value(shape_or_item: Any, key: str) -> int:
    if isinstance(shape_or_item, dict):
        return int(shape_or_item.get(key, 0) or 0)
    return int(getattr(shape_or_item, key, 0) or 0)


def center(shape_or_item: Any) -> tuple[float, float]:
    left = get_position_value(shape_or_item, "left")
    top = get_position_value(shape_or_item, "top")
    width = get_position_value(shape_or_item, "width")
    height = get_position_value(shape_or_item, "height")
    return left + width / 2, top + height / 2


def nearby_text_for_image(shape, text_items: list[dict[str, Any]], limit: int = 5) -> str:
    image_x, image_y = center(shape)
    ranked: list[tuple[float, dict[str, Any]]] = []

    for item in text_items:
        text_x, text_y = center(item)
        distance = ((image_x - text_x) ** 2 + (image_y - text_y) ** 2) ** 0.5
        ranked.append((distance, item))

    nearest = [item["text"] for _, item in sorted(ranked, key=lambda pair: pair[0])[:limit]]
    return clean_text(" | ".join(nearest))


def get_alt_text(shape) -> str:
    try:
        c_nv_pr = shape._element.xpath(".//p:cNvPr")[0]
        return clean_text(c_nv_pr.get("descr"))
    except (AttributeError, IndexError):
        return ""


def presentation_metadata(prs: Presentation) -> dict[str, Any]:
    props = prs.core_properties
    return {
        "title": json_safe(props.title),
        "subject": json_safe(props.subject),
        "author": json_safe(props.author),
        "keywords": json_safe(props.keywords),
        "comments": json_safe(props.comments),
        "category": json_safe(props.category),
        "created": json_safe(props.created),
        "modified": json_safe(props.modified),
        "last_modified_by": json_safe(props.last_modified_by),
        "slide_width_inches": emu_to_inches(prs.slide_width),
        "slide_height_inches": emu_to_inches(prs.slide_height),
    }


def extract_images_from_pptx(pptx_path: Path, output_dir: Path) -> list[dict[str, Any]]:
    prs = Presentation(str(pptx_path))
    images_dir = output_dir / "images"
    images_dir.mkdir(parents=True, exist_ok=True)

    rows: list[dict[str, Any]] = []
    for slide_number, slide in enumerate(prs.slides, start=1):
        text_items = get_slide_text_items(slide)
        slide_title = get_slide_title(slide, text_items)
        slide_text = clean_text(" | ".join(item["text"] for item in text_items))

        for image_index, shape in enumerate(iter_picture_shapes(slide.shapes), start=1):
            image = shape.image
            ext = image.ext or "bin"
            filename = (
                f"slide_{slide_number:03d}_image_{image_index:02d}_"
                f"{safe_filename(slide_title)}.{ext}"
            )
            image_path = images_dir / filename
            image_path.write_bytes(image.blob)

            row = {
                "pptx_file": str(pptx_path),
                "slide_number": slide_number,
                "image_index": image_index,
                "slide_title": slide_title,
                "image_file": str(image_path),
                "image_filename": filename,
                "image_ext": ext,
                "content_type": image.content_type,
                "image_width_px": image.size[0] if image.size else None,
                "image_height_px": image.size[1] if image.size else None,
                "image_dpi_x": image.dpi[0] if image.dpi else None,
                "image_dpi_y": image.dpi[1] if image.dpi else None,
                "shape_name": getattr(shape, "name", ""),
                "alt_text": get_alt_text(shape),
                "left_inches": emu_to_inches(getattr(shape, "left", None)),
                "top_inches": emu_to_inches(getattr(shape, "top", None)),
                "width_inches": emu_to_inches(getattr(shape, "width", None)),
                "height_inches": emu_to_inches(getattr(shape, "height", None)),
                "nearby_text": nearby_text_for_image(shape, text_items),
                "slide_text": slide_text,
            }
            rows.append(row)

    metadata_path = output_dir / "presentation_metadata.json"
    metadata_path.write_text(
        json.dumps(presentation_metadata(prs), indent=2, ensure_ascii=True),
        encoding="utf-8",
    )

    csv_path = output_dir / "image_metadata.csv"
    if rows:
        with csv_path.open("w", newline="", encoding="utf-8") as csv_file:
            writer = csv.DictWriter(csv_file, fieldnames=list(rows[0].keys()))
            writer.writeheader()
            writer.writerows(rows)
    else:
        csv_path.write_text("", encoding="utf-8")

    json_path = output_dir / "image_metadata.json"
    json_path.write_text(json.dumps(rows, indent=2, ensure_ascii=True), encoding="utf-8")

    return rows


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(
        description="Extract embedded images from a PowerPoint deck with slide title and metadata."
    )
    parser.add_argument("pptx", type=Path, help="Path to a .pptx file")
    parser.add_argument(
        "-o",
        "--output-dir",
        type=Path,
        default=None,
        help="Output folder. Defaults to <pptx name>_images.",
    )
    return parser.parse_args()


def main() -> int:
    args = parse_args()
    pptx_path = args.pptx.expanduser().resolve()

    if pptx_path.suffix.lower() == ".ppt":
        raise SystemExit("Legacy .ppt files are not supported. Save/export the file as .pptx first.")
    if pptx_path.suffix.lower() != ".pptx":
        raise SystemExit("Input must be a .pptx file.")
    if not pptx_path.exists():
        raise SystemExit(f"File not found: {pptx_path}")

    output_dir = args.output_dir or pptx_path.with_name(f"{pptx_path.stem}_images")
    output_dir = output_dir.expanduser().resolve()
    output_dir.mkdir(parents=True, exist_ok=True)

    rows = extract_images_from_pptx(pptx_path, output_dir)
    print(f"Extracted {len(rows)} image(s)")
    print(f"Images: {output_dir / 'images'}")
    print(f"Image metadata CSV: {output_dir / 'image_metadata.csv'}")
    print(f"Image metadata JSON: {output_dir / 'image_metadata.json'}")
    print(f"Presentation metadata JSON: {output_dir / 'presentation_metadata.json'}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
